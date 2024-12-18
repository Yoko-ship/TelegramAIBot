import telebot
import google.generativeai as genai
import os
from PIL import Image
import requests
from io import BytesIO
from pptx import Presentation
import mimetypes
from flask import Flask,request

app = Flask(__name__)


API_TOKEN = "7429366923:AAHGpTLn2wjz7S1jr01ttdj-_Vmz00ma3l8"
bot = telebot.TeleBot(API_TOKEN)


apiKey = "AIzaSyCrj3saz9DtSmuesXjHKLR7HIAxRJD3RrY"
genai.configure(api_key=apiKey)
model = genai.GenerativeModel("gemini-1.5-flash")
history = [
    {"role":"user","parts":"Здарова"},
    {"role":"model","parts":"Привет!"}
]


@app.route("/setup_webhook",methods=["GET","POST"])
def setup_webhook():
    webhook_url = "https://telegramaibot-9e9s.onrender.com/webhook"
    success = bot.set_webhook(url=webhook_url)
    if success:
        return "Webhook установлен", 200
    else:
        return "Ошибка установки вебхука", 500

@app.route("/webhook",methods=["POST"])
def webhook():
    try:

        json_str = request.get_data(as_text=True)
        update = telebot.types.Update.de_json(json_str)
        bot.process_new_updates([update])
        return "Ok",200
    except Exception as e:
         print(f"Ошибка обработки вебхука: {e}")


@app.route("/",methods=["GET"])
def home():
    return "Сервер работает успешно",200

def presentation_to_text(path_file):
    try:
        prs = Presentation(path_file)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception:
        return "Ошибка при извлечении текста из презентации."

@bot.message_handler(content_types=["text", "document"])
def get_text_messages(message):
    if message.text == "Привет":
        bot.send_message(message.from_user.id, "Привет чем могу помочь")
    elif message.text in ["Как дела", "Как дела?"]:
        bot.send_message(message.from_user.id, "У меня все отлично, как у тебя?")
    elif message.text == "/help":
        bot.send_message(message.from_user.id, "Доступные команды: /help, /info, /image, /pdf, /present")
    elif message.text == "/image":
        bot.send_message(message.from_user.id, "Отправьте ссылку на изображение.")
        bot.register_next_step_handler(message, get_url)
    elif message.text == "/info":
        bot.send_message(message.from_user.id, "Я тут, чтобы помочь! 🫡 Не стесняйся спрашивать.")
    elif message.text == "/present":
        bot.send_message(message.from_user.id, "Отправьте презентацию.")
        bot.register_next_step_handler(message, present_handler)
    elif message.text == "/pdf":
        bot.send_message(message.from_user.id, "Отправьте PDF файл.")
        bot.register_next_step_handler(message, pdf_handler)
    else:
        try:
            history.append({"role":"user","parts":message.text})

            chat = model.start_chat(history=history)
            response = chat.send_message(message.text)
            history.append({"role":"model","parts":response.text})
            bot.send_message(message.from_user.id, response.text.replace("*", "")[:4000])
        except Exception as e:
            bot.send_message(message.from_user.id, f"Ошибка: {e}")

def get_url(message):
    url = message.text
    if url == "/stop":
        bot.reply_to(message, "Вы приостановили функцию.")
        bot.register_next_step_handler(message, get_text_messages)
        return

    try:
        response_url = requests.get(url)
        picture_describe = Image.open(BytesIO(response_url.content))
        response_chat = model.generate_content(["Опиши эту картинку на русском", picture_describe])
        bot.send_message(message.from_user.id, response_chat.text)
    except Exception as e:
        bot.send_message(message.from_user.id, f"Ошибка: Неправильный URL или данные. Попробуйте снова. {e}")
        bot.register_next_step_handler(message, get_url)

def present_handler(message):
    if message.text == "/stop":
        bot.reply_to(message, "Вы приостановили функцию.")
        bot.register_next_step_handler(message, get_text_messages)
        return

    try:
        file_info = bot.get_file(message.document.file_id)
        downloaded_file = bot.download_file(file_info.file_path)

        file_name = message.document.file_name
        with open(file_name, "wb") as new_file:
            new_file.write(downloaded_file)

        pdfto_text = presentation_to_text(file_name)
        response = model.generate_content(["Опиши это кратко, но подробно, желательно с примерами", pdfto_text])
        bot.send_message(message.from_user.id, response.text.replace("*", "")[:4096])
    except Exception as e:
        bot.send_message(message.from_user.id, f"Ошибка: {e}")

def pdf_handler(message):
    if message.text == "/stop":
        bot.reply_to(message, "Вы приостановили функцию.")
        bot.register_next_step_handler(message, get_text_messages)
        return

    try:
        pdf_info = bot.get_file(message.document.file_id)
        downloaded_pdf = bot.download_file(pdf_info.file_path)
        pdf_name = message.document.file_name

        with open(pdf_name, "wb") as new_pdf:
            new_pdf.write(downloaded_pdf)

        mime_type, _ = mimetypes.guess_type(pdf_name)
        if mime_type is None:
            mime_type = "application/pdf"

        with open(pdf_name, "rb") as pdf_file:
            sample_pdf = genai.upload_file(pdf_file, mime_type=mime_type)

        response = model.generate_content(["Опиши это кратко, но подробно, желательно с примерами", sample_pdf])
        bot.send_message(message.from_user.id, response.text.replace("*", "")[:4096])
    except Exception as e:
        bot.reply_to(message, f"Ошибка: {e}")

if __name__ == "__main__":
    port = int(os.getenv("PORT",5000))
    print(f"Starting Flask on port {port}...")
    app.run(host='0.0.0.0', port=port)
